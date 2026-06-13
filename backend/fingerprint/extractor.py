"""
Text extraction helpers for supported file types.
"""

from pathlib import Path


def extract_text(file_path: str) -> str:
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"文件不存在: {file_path}")
    if path.stat().st_size == 0:
        raise ValueError(f"文件内容为空: {file_path}")

    ext = path.suffix.lower()
    extractors = {
        ".txt": _extract_txt,
        ".md": _extract_txt,
        ".pdf": _extract_pdf,
        ".docx": _extract_docx,
        ".pptx": _extract_pptx,
    }

    extractor = extractors.get(ext)
    if extractor is None:
        raise ValueError(
            f"不支持的文件格式: {ext}，支持: {', '.join(sorted(extractors.keys()))}"
        )

    text = extractor(file_path)
    lines = [line.strip() for line in text.splitlines()]
    return "\n".join(line for line in lines if line)


def _extract_txt(file_path: str) -> str:
    for encoding in ("utf-8", "gbk", "gb2312", "latin-1"):
        try:
            with open(file_path, "r", encoding=encoding) as file:
                return file.read()
        except (UnicodeDecodeError, LookupError):
            continue

    with open(file_path, "r", encoding="utf-8", errors="replace") as file:
        return file.read()


def _extract_pdf(file_path: str) -> str:
    try:
        from PyPDF2 import PdfReader
    except ImportError as exc:
        raise ImportError("需要安装 PyPDF2") from exc

    reader = PdfReader(file_path)
    parts = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            parts.append(text)
    return "\n".join(parts)


def _extract_docx(file_path: str) -> str:
    try:
        from docx import Document
    except ImportError as exc:
        raise ImportError("需要安装 python-docx") from exc

    document = Document(file_path)
    parts = []

    for para in document.paragraphs:
        text = para.text.strip()
        if text:
            parts.append(text)

    for table in document.tables:
        for row in table.rows:
            row_text = " | ".join(
                cell.text.strip() for cell in row.cells if cell.text.strip()
            )
            if row_text:
                parts.append(row_text)

    return "\n".join(parts)


def _extract_pptx(file_path: str) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise ImportError("需要安装 python-pptx") from exc

    presentation = Presentation(file_path)
    slides_text = []

    for slide in presentation.slides:
        texts = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    row_text = " | ".join(
                        cell.text.strip() for cell in row.cells if cell.text.strip()
                    )
                    if row_text:
                        texts.append(row_text)
        if texts:
            slides_text.append("\n".join(texts))

    return "\n".join(slides_text)
